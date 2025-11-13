"""
Generador de Presentación PowerPoint para AutoForense
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
import os

def crear_presentacion():
    """Crea la presentación de PowerPoint sobre AutoForense"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Colores del tema
    COLOR_PRINCIPAL = RGBColor(25, 118, 210)  # Azul
    COLOR_ACENTO = RGBColor(46, 125, 50)      # Verde
    COLOR_ALERTA = RGBColor(211, 47, 47)      # Rojo
    COLOR_TEXTO = RGBColor(33, 33, 33)        # Gris oscuro
    COLOR_BLANCO = RGBColor(255, 255, 255)    # Blanco
    
    # ==================== DIAPOSITIVA 1: PORTADA ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Fondo de gradiente simulado
    background = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(15, 76, 129)
    background.line.color.rgb = RGBColor(15, 76, 129)
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "AutoForense"
    p = tf.paragraphs[0]
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLANCO
    p.alignment = PP_ALIGN.CENTER
    
    # Subtítulo
    txBox = slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    tf.text = "Automatización Inteligente de Análisis Forense en Windows"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.color.rgb = COLOR_BLANCO
    p.alignment = PP_ALIGN.CENTER
    
    # Info adicional
    txBox = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "Análisis Forense Digital + Inteligencia Artificial"
    p = tf.paragraphs[0]
    p.font.size = Pt(18)
    p.font.italic = True
    p.font.color.rgb = RGBColor(200, 230, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # ==================== DIAPOSITIVA 2: ¿QUÉ ES AUTOFORENSE? ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = "¿Qué es AutoForense?"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRINCIPAL
    
    # Contenido
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    # Descripción
    p = tf.add_paragraph()
    p.text = "Herramienta de análisis forense digital para sistemas Windows que combina:"
    p.font.size = Pt(20)
    p.font.color.rgb = COLOR_TEXTO
    p.space_after = Pt(12)
    
    items = [
        ("🔍", "Recopilación automatizada de datos forenses", COLOR_ACENTO),
        ("🤖", "Análisis inteligente con IA (Google Gemini 2.5 Pro)", RGBColor(219, 68, 55)),
        ("📊", "Generación de reportes profesionales en PDF", RGBColor(251, 140, 0)),
        ("⚡", "Detección de amenazas y comportamientos sospechosos", COLOR_ALERTA),
        ("🛡️", "Operaciones no destructivas (solo lectura)", COLOR_ACENTO)
    ]
    
    for emoji, texto, color in items:
        p = tf.add_paragraph()
        p.text = f"{emoji}  {texto}"
        p.font.size = Pt(18)
        p.font.color.rgb = color
        p.level = 1
        p.space_after = Pt(10)
    
    # ==================== DIAPOSITIVA 3: ARQUITECTURA ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Arquitectura del Sistema"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRINCIPAL
    
    # Diagrama de arquitectura (texto)
    componentes = [
        ("AutoForense.py", "Interfaz Principal", Inches(1), Inches(1.5), COLOR_PRINCIPAL),
        ("PowerShell\nHelper", "Puente Python-PS", Inches(0.8), Inches(3), RGBColor(52, 152, 219)),
        ("FuncionesForenses\n.psm1", "Módulo PowerShell", Inches(0.8), Inches(4.5), RGBColor(46, 125, 50)),
        ("AIAnalyzer", "Análisis con IA", Inches(4), Inches(3), RGBColor(156, 39, 176)),
        ("Google Gemini\nAPI", "Inteligencia\nArtificial", Inches(4), Inches(4.5), RGBColor(219, 68, 55)),
        ("PDFGenerator", "Generación de\nReportes", Inches(7), Inches(3), RGBColor(255, 152, 0))
    ]
    
    for nombre, desc, left, top, color in componentes:
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, Inches(2), Inches(0.9)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(200, 200, 200)
        shape.line.width = Pt(1)
        
        tf = shape.text_frame
        tf.text = nombre
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLANCO
        p.alignment = PP_ALIGN.CENTER
        
        # Descripción
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_BLANCO
        p.alignment = PP_ALIGN.CENTER
    
    # Flechas simuladas con líneas
    flechas = [
        (Inches(2), Inches(2.1), Inches(2), Inches(2.8)),  # AutoForense -> PowerShell
        (Inches(1.8), Inches(4.95), Inches(2.8), Inches(4.95)),  # PowerShell -> Funciones
        (Inches(2), Inches(2.1), Inches(4), Inches(3.4)),  # AutoForense -> AI
        (Inches(5), Inches(3.95), Inches(5), Inches(4.5)),  # AI -> Gemini
        (Inches(3), Inches(2.1), Inches(7), Inches(3.4))  # AutoForense -> PDF
    ]
    
    for x1, y1, x2, y2 in flechas:
        connector = slide.shapes.add_connector(2, x1, y1, x2, y2)
        connector.line.color.rgb = RGBColor(100, 100, 100)
        connector.line.width = Pt(2)
    
    # ==================== DIAPOSITIVA 4: COMPONENTES PRINCIPALES ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Componentes Principales"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRINCIPAL
    
    # Columna 1
    componentes_info = [
        ("AutoForense.py", "Programa principal con menú interactivo"),
        ("FuncionesForenses.psm1", "Módulo PowerShell para recopilación"),
        ("PowershellHelper.py", "Puente Python-PowerShell"),
        ("AIAnalyzer.py", "Integración con Google Gemini"),
        ("PDFGenerator.py", "Generador de reportes PDF"),
        ("Prompt.txt", "Configuración del comportamiento de IA")
    ]
    
    top = Inches(1.2)
    for nombre, desc in componentes_info:
        # Caja de componente
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0.5), top, Inches(9), Inches(0.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(240, 240, 240)
        shape.line.color.rgb = COLOR_PRINCIPAL
        shape.line.width = Pt(2)
        
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = nombre
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRINCIPAL
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXTO
        
        top += Inches(0.95)
    
    # ==================== DIAPOSITIVA 5: FUNCIONALIDADES ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Funcionalidades"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRINCIPAL
    
    # Modo Básico
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.5), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "🔧 Modo Básico (sin IA)"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    
    txBox = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(4.3), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    funciones_basicas = [
        "Get-SuspiciousEvents\nExtrae eventos sospechosos del Visor de Eventos",
        "Get-InternetProcesses\nCorrelaciona procesos con conexiones de red",
        "Get-UnsignedProcesses\nDetecta procesos sin firma digital"
    ]
    
    for func in funciones_basicas:
        p = tf.add_paragraph()
        p.text = "• " + func
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(12)
    
    # Modo con IA
    txBox = slide.shapes.add_textbox(Inches(5.5), Inches(1.2), Inches(4), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "🤖 Modo con IA"
    p = tf.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(156, 39, 176)
    
    txBox = slide.shapes.add_textbox(Inches(5.7), Inches(1.9), Inches(3.8), Inches(2.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    funciones_ia = [
        "Análisis Forense con IA\nEjecuta tarea específica + análisis IA",
        "Análisis Forense Completo\nEjecuta todas las tareas + reporte consolidado"
    ]
    
    for func in funciones_ia:
        p = tf.add_paragraph()
        p.text = "• " + func
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(12)
    
    # Características adicionales
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.5), Inches(4.8), Inches(9), Inches(2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 248, 225)
    shape.line.color.rgb = RGBColor(255, 152, 0)
    shape.line.width = Pt(2)
    
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "✨ Características Destacadas"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(230, 81, 0)
    p.space_after = Pt(10)
    
    caracteristicas = [
        "✓ Detección automática de amenazas y comportamientos anómalos",
        "✓ Reportes profesionales en PDF con análisis detallado",
        "✓ Correlación inteligente de datos forenses",
        "✓ Recomendaciones de seguridad personalizadas"
    ]
    
    for carac in caracteristicas:
        p = tf.add_paragraph()
        p.text = carac
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(6)
    
    # ==================== DIAPOSITIVA 6: FLUJO DE TRABAJO ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Flujo de Trabajo"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRINCIPAL
    
    # Pasos del flujo
    pasos = [
        ("1", "Usuario selecciona tareas", RGBColor(33, 150, 243)),
        ("2", "PowerShell recopila datos del sistema", RGBColor(76, 175, 80)),
        ("3", "Datos se exportan en formato CSV", RGBColor(255, 152, 0)),
        ("4", "IA analiza y detecta anomalías", RGBColor(156, 39, 176)),
        ("5", "Se genera reporte PDF profesional", RGBColor(244, 67, 54))
    ]
    
    top = Inches(1.3)
    for num, desc, color in pasos:
        # Círculo con número
        circle = slide.shapes.add_shape(
            3,  # Oval
            Inches(0.8), top, Inches(0.6), Inches(0.6)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = color
        circle.line.color.rgb = color
        
        tf = circle.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = num
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLANCO
        p.alignment = PP_ALIGN.CENTER
        
        # Descripción
        txBox = slide.shapes.add_textbox(Inches(1.7), top, Inches(7.5), Inches(0.6))
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_TEXTO
        
        # Flecha hacia abajo (excepto el último)
        if num != "5":
            connector = slide.shapes.add_connector(
                2,  # Straight connector
                Inches(1.1), top + Inches(0.7), Inches(1.1), top + Inches(0.95)
            )
            connector.line.color.rgb = RGBColor(150, 150, 150)
            connector.line.width = Pt(3)
        
        top += Inches(1.05)
    
    # ==================== DIAPOSITIVA 7: INSTALACIÓN ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Instalación y Configuración"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRINCIPAL
    
    # Requisitos
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(8.4), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "📋 Requisitos"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    
    p = tf.add_paragraph()
    p.text = "• Windows 10/11  • Python 3.8+  • PowerShell 5.1+  • API Key de Google AI"
    p.font.size = Pt(16)
    p.font.color.rgb = COLOR_TEXTO
    p.space_before = Pt(8)
    
    # Pasos de instalación
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(2.5), Inches(8.4), Inches(4))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "⚙️ Pasos de Instalación"
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    p.space_after = Pt(12)
    
    pasos_instalacion = [
        ("1.", "Clonar el repositorio", "git clone https://github.com/usuario/PIA-PC.git"),
        ("2.", "Instalar dependencias Python", "pip install -r requirements.txt"),
        ("3.", "Configurar API Key de Google AI", "Crear archivo .env con: GOOGLE_API_KEY=tu_key"),
        ("4.", "Ejecutar el programa", ".\\ejecutar_autoforense.bat")
    ]
    
    for num, titulo, comando in pasos_instalacion:
        p = tf.add_paragraph()
        p.text = f"{num} {titulo}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRINCIPAL
        p.space_after = Pt(4)
        
        p = tf.add_paragraph()
        p.text = comando
        p.font.size = Pt(13)
        p.font.name = "Consolas"
        p.font.color.rgb = RGBColor(100, 100, 100)
        p.level = 1
        p.space_after = Pt(10)
    
    # ==================== DIAPOSITIVA 8: DEMO/USO ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Modo de Uso"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRINCIPAL
    
    # Menú principal simulado
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(30, 30, 30)
    shape.line.color.rgb = RGBColor(100, 100, 100)
    shape.line.width = Pt(1)
    
    tf = shape.text_frame
    tf.word_wrap = True
    
    menu_text = """AutoForense - Menú Principal

Funciones disponibles:

1. Get-SuspiciousEvents - Extraer eventos sospechosos
2. Get-InternetProcesses - Procesos con conexiones de red
3. Get-UnsignedProcesses - Procesos sin firma digital
4. Análisis Forense con IA - Ejecuta tarea y analiza con IA
5. Análisis Forense Completo - Todas las tareas + reporte PDF
6. Salir

Seleccione una opción (1-6): _"""
    
    p = tf.paragraphs[0]
    p.text = menu_text
    p.font.size = Pt(16)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(0, 255, 0)
    
    # ==================== DIAPOSITIVA 9: TECNOLOGÍAS ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Tecnologías Utilizadas"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRINCIPAL
    
    # Tecnologías en cajas
    tecnologias = [
        ("Python 3.8+", "Lenguaje principal", RGBColor(53, 114, 165)),
        ("PowerShell 5.1+", "Recopilación forense", RGBColor(1, 36, 86)),
        ("Google Gemini 2.5", "Inteligencia Artificial", RGBColor(219, 68, 55)),
        ("ReportLab", "Generación de PDF", RGBColor(244, 67, 54)),
        ("python-pptx", "Presentaciones", RGBColor(209, 52, 56)),
        ("Windows API", "Análisis del sistema", RGBColor(0, 120, 215))
    ]
    
    # Distribución en 2 columnas
    positions = [
        (Inches(0.8), Inches(1.5)),
        (Inches(5.2), Inches(1.5)),
        (Inches(0.8), Inches(3)),
        (Inches(5.2), Inches(3)),
        (Inches(0.8), Inches(4.5)),
        (Inches(5.2), Inches(4.5))
    ]
    
    for (nombre, desc, color), (left, top) in zip(tecnologias, positions):
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            left, top, Inches(3.8), Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.color.rgb = RGBColor(200, 200, 200)
        shape.line.width = Pt(2)
        
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = nombre
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = COLOR_BLANCO
        p.alignment = PP_ALIGN.CENTER
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_BLANCO
        p.alignment = PP_ALIGN.CENTER
    
    # ==================== DIAPOSITIVA 10: RESULTADOS ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Resultados y Reportes"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRINCIPAL
    
    # Tipos de salida
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "📊 Tipos de Salida"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    p.space_after = Pt(15)
    
    resultados = [
        ("Archivos CSV", "Datos tabulados de eventos, procesos y conexiones"),
        ("Reportes PDF", "Análisis completo con hallazgos y recomendaciones"),
        ("Consola", "Resultados inmediatos en tiempo real")
    ]
    
    for tipo, desc in resultados:
        p = tf.add_paragraph()
        p.text = f"• {tipo}"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = COLOR_PRINCIPAL
        p.space_after = Pt(6)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(16)
        p.font.color.rgb = COLOR_TEXTO
        p.level = 1
        p.space_after = Pt(12)
    
    # Información del reporte PDF
    shape = slide.shapes.add_shape(
        1,  # Rectangle
        Inches(0.8), Inches(4.2), Inches(8.4), Inches(2.3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(232, 245, 233)
    shape.line.color.rgb = COLOR_ACENTO
    shape.line.width = Pt(2)
    
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = "📄 Contenido del Reporte PDF"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    p.space_after = Pt(10)
    
    contenido_pdf = [
        "✓ Resumen ejecutivo de hallazgos",
        "✓ Análisis detallado de amenazas detectadas",
        "✓ Estadísticas y métricas del sistema",
        "✓ Recomendaciones de seguridad específicas",
        "✓ Nivel de riesgo por categoría"
    ]
    
    for item in contenido_pdf:
        p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = COLOR_TEXTO
        p.space_after = Pt(5)
    
    # ==================== DIAPOSITIVA 11: SEGURIDAD ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    tf.text = "Consideraciones de Seguridad"
    p = tf.paragraphs[0]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRINCIPAL
    
    # Principios de seguridad
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(8.4), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    p = tf.add_paragraph()
    p.text = "🛡️ Principios de Seguridad"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = COLOR_ACENTO
    p.space_after = Pt(15)
    
    principios = [
        ("No Destructivo", "Solo operaciones de lectura, sin modificar el sistema"),
        ("Autorización", "Obtener permisos antes de analizar sistemas"),
        ("Privacidad", "No compartir reportes con información sensible"),
        ("Verificación", "Los hallazgos deben ser validados por profesionales")
    ]
    
    for titulo, desc in principios:
        shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(1), Inches(1.8 + (principios.index((titulo, desc)) * 1.05)), 
            Inches(8), Inches(0.85)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 243, 224)
        shape.line.color.rgb = RGBColor(255, 152, 0)
        shape.line.width = Pt(1.5)
        
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = f"🔒 {titulo}"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = RGBColor(230, 81, 0)
        
        p = tf.add_paragraph()
        p.text = desc
        p.font.size = Pt(14)
        p.font.color.rgb = COLOR_TEXTO
    
    # Nota legal
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(6.1), Inches(8.4), Inches(1))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "⚖️ AutoForense se proporciona 'tal cual', sin garantías. El usuario es responsable del cumplimiento de leyes y regulaciones aplicables."
    p.font.size = Pt(12)
    p.font.italic = True
    p.font.color.rgb = RGBColor(100, 100, 100)
    p.alignment = PP_ALIGN.CENTER
    
    # ==================== DIAPOSITIVA 12: CONCLUSIONES ====================
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Fondo similar a portada
    background = slide.shapes.add_shape(
        1,  # Rectangle
        0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(15, 76, 129)
    background.line.color.rgb = RGBColor(15, 76, 129)
    
    # Título
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(8), Inches(1))
    tf = txBox.text_frame
    tf.text = "AutoForense"
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLANCO
    p.alignment = PP_ALIGN.CENTER
    
    # Conclusiones
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.8), Inches(7), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    
    conclusiones = [
        "✅ Análisis forense automatizado e inteligente",
        "✅ Detección temprana de amenazas",
        "✅ Reportes profesionales accionables",
        "✅ Fácil de usar, potente en resultados",
        "✅ Código abierto y extensible"
    ]
    
    for conclusion in conclusiones:
        p = tf.add_paragraph()
        p.text = conclusion
        p.font.size = Pt(20)
        p.font.color.rgb = COLOR_BLANCO
        p.space_after = Pt(12)
    
    # Gracias
    txBox = slide.shapes.add_textbox(Inches(1), Inches(5.8), Inches(8), Inches(0.8))
    tf = txBox.text_frame
    tf.text = "¡Gracias por su atención!"
    p = tf.paragraphs[0]
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLANCO
    p.alignment = PP_ALIGN.CENTER
    
    # Información de contacto/repositorio
    txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    tf = txBox.text_frame
    tf.text = "GitHub: github.com/usuario/PIA-PC"
    p = tf.paragraphs[0]
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(200, 230, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Guardar presentación
    output_file = "AutoForense_Presentacion.pptx"
    prs.save(output_file)
    print(f"\n✓ Presentación creada exitosamente: {output_file}")
    print(f"✓ Total de diapositivas: {len(prs.slides)}")
    print(f"✓ Ubicación: {os.path.abspath(output_file)}\n")
    
    return output_file

if __name__ == "__main__":
    print("\n" + "="*60)
    print("   GENERADOR DE PRESENTACIÓN - AUTOFORENSE")
    print("="*60)
    
    try:
        # Verificar si python-pptx está instalado
        try:
            from pptx import Presentation
        except ImportError:
            print("\n⚠ python-pptx no está instalado")
            print("Instalando dependencias...")
            import subprocess
            import sys
            subprocess.check_call([sys.executable, '-m', 'pip', 'install', 'python-pptx'])
            print("✓ Dependencias instaladas\n")
        
        # Crear presentación
        archivo = crear_presentacion()
        
        print("="*60)
        print("La presentación está lista para usarse.")
        print("="*60 + "\n")
        
    except Exception as e:
        print(f"\n✗ Error al crear la presentación: {e}\n")
        import traceback
        traceback.print_exc()

